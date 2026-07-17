#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""生成宠寻寻 - AI图像识别寻宠平台 答辩PPT (V2 - 修改版)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 颜色主题 ──
COLOR_PRIMARY = RGBColor(0x66, 0x7E, 0xEA)
COLOR_SECONDARY = RGBColor(0x76, 0x4B, 0xA2)
COLOR_ACCENT = RGBColor(0xFF, 0x6B, 0x6B)
COLOR_DARK = RGBColor(0x2D, 0x2D, 0x2D)
COLOR_GRAY = RGBColor(0x66, 0x66, 0x66)
COLOR_LIGHT = RGBColor(0xF8, 0xF9, 0xFA)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_GREEN = RGBColor(0x4C, 0xAF, 0x50)
COLOR_ORANGE = RGBColor(0xFF, 0x98, 0x00)
COLOR_YELLOW = RGBColor(0xFF, 0xC1, 0x07)
COLOR_BAR_BG = RGBColor(0xE8, 0xEC, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_background(slide, color=COLOR_WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, color=COLOR_DARK, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_bullet_textbox(slide, left, top, width, height, items, font_size=16, color=COLOR_DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return txBox


def add_page_header(slide, title, subtitle=None):
    add_shape(slide, 0, 0, W, Inches(1.1), COLOR_PRIMARY)
    add_textbox(slide, Inches(0.8), Inches(0.12), Inches(10), Inches(0.65), title,
                font_size=30, color=COLOR_WHITE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(0.65), Inches(10), Inches(0.38), subtitle,
                    font_size=14, color=RGBColor(0xE0, 0xE0, 0xE0))
    add_shape(slide, 0, Inches(1.1), W, Inches(0.06), COLOR_SECONDARY)


def add_page_number(slide, num, total):
    add_textbox(slide, Inches(12), Inches(7.0), Inches(1.2), Inches(0.4),
                f"{num}/{total}", font_size=12, color=COLOR_GRAY, alignment=PP_ALIGN.RIGHT)


def add_bar_chart(slide, left, top, width, data, max_val=5.0, bar_height=0.38, color=None):
    """绘制水平条形图 data=[(label, value, color), ...]"""
    if color is None:
        color = COLOR_PRIMARY
    row_h = bar_height + 0.12
    bar_max_w = width - Inches(4.2)
    for i, (label, val, bar_color) in enumerate(data):
        y = top + i * row_h
        # 标签
        add_textbox(slide, left, y, Inches(3.0), Inches(bar_height),
                    f"  {label}", font_size=12, color=COLOR_DARK)
        # 背景条
        add_shape(slide, left + Inches(3.0), y + Inches(0.04),
                  bar_max_w, Inches(bar_height - 0.04), COLOR_BAR_BG)
        # 数值条
        fill_w = bar_max_w * (val / max_val)
        add_shape(slide, left + Inches(3.0), y + Inches(0.04),
                  fill_w, Inches(bar_height - 0.04), bar_color)
        # 数值
        add_textbox(slide, left + Inches(3.1) + fill_w, y,
                    Inches(1.0), Inches(bar_height),
                    f" {val:.1f}", font_size=11, color=bar_color, bold=True)


TOTAL_SLIDES = 14

# ══════════════════════════════════════════════
# Slide 1: 封面
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_PRIMARY)
add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
            "宠寻寻", font_size=54, color=COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
            "AI图像识别寻宠平台", font_size=36, color=RGBColor(0xE0, 0xE0, 0xFF), alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(4.5), Inches(3.8), Inches(4), Inches(0.04), COLOR_WHITE)
add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
            "设计思维期末项目答辩", font_size=28, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
            "团队成员：产品经理 · 技术开发 · UI设计", font_size=18, color=RGBColor(0xCC, 0xCC, 0xFF), alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.5),
            "答辩日期：2026年6月23日", font_size=16, color=RGBColor(0xCC, 0xCC, 0xFF), alignment=PP_ALIGN.CENTER)
add_page_number(slide, 1, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# Slide 2: 目录
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_WHITE)
add_page_header(slide, "目  录", "CONTENTS")
items = [
    "1.  项目背景与问题发现",
    "2.  用户需求分析",
    "3.  产品解决方案",
    "4.  产品原型展示",
    "5.  Demo 演示",
    "6.  商业模式设计",
    "7.  市场分析与竞争分析",
    "8.  技术架构介绍",
    "9.  团队分工",
    "10. 项目总结与未来规划"
]
txBox = add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(5.5), "", font_size=22, color=COLOR_DARK)
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(items):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(22)
    p.font.color.rgb = COLOR_DARK
    p.space_after = Pt(10)
add_page_number(slide, 2, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# Slide 3: 项目背景与问题发现（V2 - 新增故事+数据图表）
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_WHITE)
add_page_header(slide, "一、项目背景与问题发现", "宠物走失——一个被忽视的社会问题")

# ── 故事引入（顶部暖色卡片） ──
add_shape(slide, Inches(0.6), Inches(1.35), Inches(12.1), Inches(1.25), RGBColor(0xFF, 0xF3, 0xF0))
add_textbox(slide, Inches(0.9), Inches(1.4), Inches(1.6), Inches(0.4),
            "📖 我们的故事", font_size=14, color=COLOR_ACCENT, bold=True)
add_textbox(slide, Inches(0.9), Inches(1.78), Inches(11.6), Inches(0.75),
            '团队成员的舍友养了一只橘猫"小白"，某天晚上不小心跑出了宿舍门。当我们焦急万分、四处张贴启事、人肉在校园里搜寻的时候，舍友整夜难眠、情绪崩溃——幸运的是，我们在48小时内通过朋友圈接力找到了它。',
            font_size=12, color=RGBColor(0x55, 0x33, 0x33))
add_textbox(slide, Inches(9.5), Inches(1.42), Inches(3.0), Inches(0.4),
            '正是这次经历，让我们意识到——宠物寻回是一个真实存在且日益严峻的社会需求。', font_size=11, color=RGBColor(0x88, 0x44, 0x44), bold=True)

# ── 左侧：基础数据 ──
add_textbox(slide, Inches(0.6), Inches(2.78), Inches(4.5), Inches(0.4),
            "📊 宠物走失问题现状", font_size=18, color=COLOR_PRIMARY, bold=True)
data_left = [
    "🐾 年走失宠物：超 1000 万只",
    "📉 传统寻回率：仅约 20%",
    "⏰ 48小时内寻回概率：80% 以上",
    "💔 近40%的养宠人经历过宠物走失",
]
add_bullet_textbox(slide, Inches(0.6), Inches(3.28), Inches(4.8), Inches(2.0), data_left, font_size=14)

# ── 右侧：痛点评分条形图 ──
add_textbox(slide, Inches(5.7), Inches(2.78), Inches(6.5), Inches(0.4),
            "⚠️ 传统寻宠方式痛点评分（1-5分）", font_size=18, color=COLOR_ACCENT, bold=True)

pain_data = [
    ("信息传播范围有限", 4.6, COLOR_ACCENT),
    ("时效性差，错过最佳时机", 4.4, RGBColor(0xFF, 0x85, 0x51)),
    ("照片难以有效比对", 4.3, RGBColor(0xFF, 0xA7, 0x26)),
    ("信息分散、缺少统一平台", 4.2, COLOR_ORANGE),
    ("悬赏机制缺乏信任", 3.8, COLOR_YELLOW),
    ("缺乏指引，不知从何做起", 3.7, RGBColor(0xA5, 0xD6, 0xA7)),
]
add_bar_chart(slide, Inches(5.7), Inches(3.28), Inches(6.8), pain_data, max_val=5.0, bar_height=0.4)

# ── 底部洞察 ──
add_shape(slide, Inches(0.6), Inches(5.8), Inches(12.1), Inches(0.9), RGBColor(0xEE, 0xF0, 0xFF))
add_textbox(slide, Inches(1.0), Inches(5.9), Inches(2.0), Inches(0.35),
            "💡 核心洞察", font_size=16, color=COLOR_PRIMARY, bold=True)
add_textbox(slide, Inches(1.0), Inches(6.2), Inches(11.3), Inches(0.45),
            "「每一只走失的宠物背后，都有一个焦急等待的家庭」—— 我们希望通过AI技术，让寻宠更高效、让失主不再无助",
            font_size=13, color=COLOR_DARK)

add_page_number(slide, 3, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# Slide 4: 用户需求分析
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_WHITE)
add_page_header(slide, "二、用户需求分析", "目标用户群体与核心需求")

card_data = [
    ("🐱 宠物主人", "快速发布走失信息\n高效匹配线索\n实时追踪寻宠进度", COLOR_PRIMARY),
    ("💚 爱心人士", "方便上报发现线索\n查看附近走失信息\n参与爱心寻宠", COLOR_GREEN),
    ("🏥 宠物机构", "批量管理走失/领养信息\n协助扩散寻宠消息\n公益合作", COLOR_ORANGE),
]
for i, (title, desc, color) in enumerate(card_data):
    x = Inches(0.8 + i * 4.1)
    y = Inches(1.8)
    add_shape(slide, x, y, Inches(3.7), Inches(3.0), RGBColor(0xF8, 0xF9, 0xFA))
    add_shape(slide, x, y, Inches(3.7), Inches(0.08), color)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.1), Inches(0.5), title, font_size=20, color=color, bold=True)
    add_textbox(slide, x + Inches(0.3), y + Inches(1.0), Inches(3.1), Inches(1.8), desc, font_size=15, color=COLOR_GRAY)

add_shape(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.2), RGBColor(0xF0, 0xF4, 0xFF))
add_textbox(slide, Inches(1.0), Inches(5.3), Inches(3), Inches(0.4), "🚶 用户旅程地图：", font_size=17, color=COLOR_PRIMARY, bold=True)
add_textbox(slide, Inches(1.0), Inches(5.8), Inches(11), Inches(0.5),
            "发现走失 → 发布信息 → 等待线索 → AI智能匹配 → 确认找回 → 安全回家", font_size=18, color=COLOR_DARK)
add_page_number(slide, 4, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# Slide 5: 产品解决方案
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_WHITE)
add_page_header(slide, "三、产品解决方案", "「宠寻寻」— 基于AI图像识别的智能寻宠平台")

add_shape(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(1.0), RGBColor(0xEE, 0xF0, 0xFF))
add_textbox(slide, Inches(1.2), Inches(1.7), Inches(11), Inches(0.4), "🎯 产品定位", font_size=20, color=COLOR_PRIMARY, bold=True)
add_textbox(slide, Inches(1.2), Inches(2.1), Inches(11), Inches(0.4),
            "面向宠物主人的AI图像识别寻宠平台，使命：让每一只走失的宠物都能安全回家", font_size=16, color=COLOR_DARK)

modules = [
    ("📝 用户管理", "手机号注册/登录\n用户信息管理\n联系方式自动同步"),
    ("🐕 宠物档案", "添加/编辑宠物信息\n上传多角度照片\n记录特殊特征标记"),
    ("🔍 走失发布", "上传走失照片\n填写走失时间、地点\n设置悬赏金额"),
    ("📸 线索上报", "上传发现照片\n标记发现地点\n描述发现情况"),
    ("🤖 AI智能匹配", "图像特征提取\n智能相似度匹配\n按匹配度排序展示"),
]
for i, (title, desc) in enumerate(modules):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(3.0 + row * 2.0)
    add_shape(slide, x, y, Inches(3.7), Inches(1.7), RGBColor(0xF8, 0xF9, 0xFA))
    add_textbox(slide, x + Inches(0.2), y + Inches(0.15), Inches(3.3), Inches(0.4), title, font_size=17, color=COLOR_PRIMARY, bold=True)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.6), Inches(3.3), Inches(1.0), desc, font_size=13, color=COLOR_GRAY)
add_page_number(slide, 5, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# Slide 6: 产品原型展示
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_WHITE)
add_page_header(slide, "四、产品原型展示", "核心页面与交互流程")

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.5), "📐 页面架构", font_size=20, color=COLOR_PRIMARY, bold=True)
arch_text = """首页
├── 走失信息列表
├── 附近发现线索
├── AI匹配入口
└── 我的（用户中心）
    ├── 宠物档案
    ├── 发布记录
    ├── 线索管理
    └── 账户设置"""
add_textbox(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(3.5), arch_text, font_size=15, color=COLOR_DARK)

add_textbox(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.5), "📱 核心页面设计", font_size=20, color=COLOR_PRIMARY, bold=True)
pages = [
    "🏠 首页 — 走失信息列表 + 地图视图切换",
    "📝 发布页 — 图片上传 + 表单填写",
    "🤖 匹配页 — AI匹配结果 + 相似度排序",
    "📍 线索页 — 线索详情 + 地图定位",
    "👤 个人中心 — 宠物档案 + 发布记录",
]
add_bullet_textbox(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(3.0), pages, font_size=15)

add_shape(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.0), RGBColor(0xEE, 0xF0, 0xFF))
add_textbox(slide, Inches(1.2), Inches(5.9), Inches(11), Inches(0.4), "🔄 核心交互流程", font_size=17, color=COLOR_PRIMARY, bold=True)
add_textbox(slide, Inches(1.2), Inches(6.3), Inches(11), Inches(0.4),
            "发布走失 → 系统推送 → 爱心人士上报线索 → AI图像匹配 → 匹配结果通知 → 联系确认 → 成功找回", font_size=15, color=COLOR_DARK)
add_page_number(slide, 6, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# Slide 7: Demo演示
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_WHITE)
add_page_header(slide, "五、Demo 演示", "核心功能完整演示流程")

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.5), "🎬 演示流程（3分钟）", font_size=22, color=COLOR_PRIMARY, bold=True)
demo_flow = [
    "① 0:00-0:20  开场介绍 — 产品定位与核心价值",
    "② 0:20-0:50  用户注册登录 — 手机号快速注册",
    "③ 0:50-1:20  创建宠物档案 — 上传照片、填写信息",
    "④ 1:20-1:50  发布走失信息 — 选择宠物、填写走失详情",
    "⑤ 1:50-2:20  上报发现线索 — 爱心人士上传发现照片",
    "⑥ 2:20-2:50  AI匹配演示 — 双重AI识别、相似度排序",
    "⑦ 2:50-3:00  总结 — 产品价值与未来规划",
]
add_bullet_textbox(slide, Inches(0.8), Inches(2.3), Inches(7), Inches(4.0), demo_flow, font_size=16)

add_shape(slide, Inches(8.5), Inches(1.6), Inches(4.2), Inches(5.0), RGBColor(0xF8, 0xF9, 0xFA))
add_textbox(slide, Inches(8.8), Inches(1.8), Inches(3.6), Inches(0.5), "✨ Demo 亮点", font_size=20, color=COLOR_PRIMARY, bold=True)
highlights = [
    "✅ 完整闭环流程可演示",
    "✅ 真实数据支撑",
    "✅ AI双重识别技术",
    "✅ 简洁易用的交互界面",
    "✅ 响应式设计适配多端",
]
add_bullet_textbox(slide, Inches(8.8), Inches(2.5), Inches(3.6), Inches(3.0), highlights, font_size=15, color=COLOR_GREEN)
add_page_number(slide, 7, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# Slide 8: 商业模式设计
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_WHITE)
add_page_header(slide, "六、商业模式设计", "可持续的商业模式探索")

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.5), "💰 收入模式", font_size=22, color=COLOR_PRIMARY, bold=True)
revenue = [
    "🔹 悬赏分成：成功寻回后收取10%赏金分成",
    "🔹 增值服务：VIP会员、优先匹配、广告推广",
    "🔹 企业合作：宠物医院、宠物店入驻费用",
    "🔹 数据服务：宠物走失数据分析报告",
]
add_bullet_textbox(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(2.5), revenue, font_size=16)

add_textbox(slide, Inches(7), Inches(1.6), Inches(5), Inches(0.5), "📊 成本结构", font_size=22, color=COLOR_PRIMARY, bold=True)
costs = [
    "🔸 技术成本：AI API调用费用 + 服务器费用",
    "🔸 运营成本：客服、推广、维护",
    "🔸 人力成本：开发团队、运营团队",
]
add_bullet_textbox(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(2.0), costs, font_size=16)

add_shape(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.8), RGBColor(0xEE, 0xF0, 0xFF))
add_textbox(slide, Inches(1.2), Inches(4.9), Inches(3), Inches(0.4), "📈 市场切入策略", font_size=18, color=COLOR_PRIMARY, bold=True)
strategy = [
    "1️⃣ 种子用户：本地宠物社群、宠物医院合作推广",
    "2️⃣ 区域扩张：从一线城市逐步向全国扩展",
    "3️⃣ 生态建设：与宠物相关企业建立长期合作",
]
add_bullet_textbox(slide, Inches(1.2), Inches(5.3), Inches(11), Inches(1.2), strategy, font_size=15)
add_page_number(slide, 8, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# Slide 9: 市场分析与竞争分析
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_WHITE)
add_page_header(slide, "七、市场分析与竞争分析", "宠物经济蓝海市场")

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.5), "📈 市场规模", font_size=22, color=COLOR_PRIMARY, bold=True)
market = [
    "🐾 宠物市场规模：超 3000 亿元",
    "🏠 养宠家庭数量：超 1 亿户",
    "🐕 年走失宠物：超 1000 万只",
    "🎯 目标市场潜力巨大",
]
add_bullet_textbox(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(2.5), market, font_size=16)

add_textbox(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.5), "🏆 竞争格局", font_size=22, color=COLOR_PRIMARY, bold=True)
comp = [
    "🔴 传统寻宠平台：用户基数大，但技术落后",
    "🟡 宠物社交APP：社交属性强，寻宠非核心",
    "🟢 本地社群：本地化强，但信息碎片化",
]
add_bullet_textbox(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(2.0), comp, font_size=15)

add_shape(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.0), RGBColor(0xE8, 0xF5, 0xE9))
add_textbox(slide, Inches(1.2), Inches(4.9), Inches(3), Inches(0.4), "✅ 我们的竞争优势", font_size=20, color=COLOR_GREEN, bold=True)
advantages = [
    "🤖 AI技术领先：双重AI识别（百度AI + 大模型），匹配准确率高",
    "📱 用户体验好：简洁易用的小程序设计，操作门槛低",
    "⚡ 实时性强：即时推送匹配结果，不错过最佳寻回时机",
    "💚 社区互助：构建爱心寻宠社区，汇聚社会力量",
]
add_bullet_textbox(slide, Inches(1.2), Inches(5.3), Inches(11), Inches(1.5), advantages, font_size=15, color=COLOR_DARK)
add_page_number(slide, 9, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# Slide 10: 技术架构介绍
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_WHITE)
add_page_header(slide, "八、技术架构介绍", "全栈技术方案")

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.5), "🛠️ 技术栈", font_size=22, color=COLOR_PRIMARY, bold=True)
tech = [
    ("前端", "HTML5 + CSS3 + Vue.js 2.x", COLOR_PRIMARY),
    ("后端", "Python Flask + RESTful API", COLOR_SECONDARY),
    ("数据库", "SQLite（轻量文件型数据库）", COLOR_GREEN),
    ("AI识别", "百度AI动物识别 + GPT-4 Vision", COLOR_ORANGE),
    ("地图", "微信小程序地图组件", COLOR_ACCENT),
]
for i, (layer, desc, color) in enumerate(tech):
    y = Inches(2.2 + i * 0.55)
    add_shape(slide, Inches(0.8), y, Inches(1.5), Inches(0.4), color)
    add_textbox(slide, Inches(0.9), y + Inches(0.02), Inches(1.3), Inches(0.35), layer, font_size=13, color=COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.5), y + Inches(0.02), Inches(4), Inches(0.35), desc, font_size=14, color=COLOR_DARK)

add_textbox(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.5), "🏗️ 系统架构", font_size=22, color=COLOR_PRIMARY, bold=True)
arch = """┌─────────────────────┐
│  用户端（Web/小程序）  │
└──────────┬──────────┘
           ↓
┌─────────────────────┐
│  API Gateway (Flask) │
└──────────┬──────────┘
           ↓
┌─────────────────────┐
│     业务服务层        │
│ 用户/宠物/走失/线索管理 │
└──────────┬──────────┘
           ↓
┌─────────────────────┐
│     AI 服务层         │
│ 百度AI + GPT-4 Vision│
└──────────┬──────────┘
           ↓
┌─────────────────────┐
│   数据层 (SQLite)    │
└─────────────────────┘"""
add_textbox(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(4.5), arch, font_size=12, color=COLOR_DARK)

add_shape(slide, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.5), RGBColor(0xEE, 0xF0, 0xFF))
add_textbox(slide, Inches(1.2), Inches(5.4), Inches(5), Inches(0.4), "🧮 AI匹配算法", font_size=18, color=COLOR_PRIMARY, bold=True)
add_textbox(slide, Inches(1.2), Inches(5.8), Inches(11), Inches(0.9),
            "相似度 = 图像相似度 + 百度AI加成(物种匹配+0.2×置信度) + 大模型特征匹配(毛色+0.1, 特征+0.05×匹配数量)", font_size=14, color=COLOR_DARK)
add_page_number(slide, 10, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# Slide 11: 团队分工（V2 - 去掉7天速通，重新排版）
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_WHITE)
add_page_header(slide, "九、团队分工", "三人协作 · 高效推进")

# ── 三人角色卡片（重新排版：左侧+右侧各一列+底部跨两列） ──
roles = [
    ("📋 产品经理", "成员A", COLOR_PRIMARY,
     ["需求调研与分析", "产品需求文档（PRD）", "商业模式设计", "答辩PPT制作与统筹", "竞品分析报告"]),
    ("💻 技术开发", "成员B", COLOR_GREEN,
     ["后端API开发（Flask + SQLite）", "AI模型集成与调优", "数据库设计与实现", "前后端联调与部署", "Demo稳定性保障"]),
    ("🎨 UI/UX设计", "成员C", COLOR_ORANGE,
     ["界面视觉设计（Figma）", "交互原型设计", "设计规范制定", "演示视频制作与剪辑", "用户体验持续优化"]),
]

# 前两个卡片：左侧并排
for i, (role, name, color, tasks) in enumerate(roles[:2]):
    x = Inches(0.8 + i * 6.2)
    y = Inches(1.5)
    add_shape(slide, x, y, Inches(5.8), Inches(2.6), RGBColor(0xF8, 0xF9, 0xFA))
    add_shape(slide, x, y, Inches(5.8), Inches(0.08), color)
    add_textbox(slide, x + Inches(0.25), y + Inches(0.2), Inches(3.0), Inches(0.4), role, font_size=20, color=color, bold=True)
    add_textbox(slide, x + Inches(3.3), y + Inches(0.22), Inches(2.2), Inches(0.35), name, font_size=13, color=COLOR_GRAY, alignment=PP_ALIGN.RIGHT)
    for j, task in enumerate(tasks):
        add_textbox(slide, x + Inches(0.25), y + Inches(0.72 + j * 0.36), Inches(5.3), Inches(0.33),
                    f"▸ {task}", font_size=12, color=COLOR_DARK)

# 第三个卡片：跨两列居中
role, name, color, tasks = roles[2]
add_shape(slide, Inches(0.8), Inches(4.25), Inches(11.7), Inches(2.0), RGBColor(0xF8, 0xF9, 0xFA))
add_shape(slide, Inches(0.8), Inches(4.25), Inches(11.7), Inches(0.08), color)
add_textbox(slide, Inches(1.05), Inches(4.43), Inches(3.0), Inches(0.4), role, font_size=20, color=color, bold=True)
add_textbox(slide, Inches(11.0), Inches(4.45), Inches(1.3), Inches(0.35), name, font_size=13, color=COLOR_GRAY, alignment=PP_ALIGN.RIGHT)
# 横排展示5个任务
for j, task in enumerate(tasks):
    x = Inches(0.8 + j * 2.35)
    add_textbox(slide, x + Inches(0.25), Inches(4.9), Inches(2.1), Inches(1.2),
                f"▸ {task}", font_size=11, color=COLOR_DARK)

add_page_number(slide, 11, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# Slide 12: 项目总结与未来规划
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_WHITE)
add_page_header(slide, "十、项目总结与未来规划", "回顾成果 · 展望未来")

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.5), "✅ 已完成工作", font_size=22, color=COLOR_GREEN, bold=True)
done = [
    "✔ 产品需求分析与用户调研",
    "✔ 后端API开发完成（Flask + SQLite）",
    "✔ AI图像识别集成（百度AI + 大模型）",
    "✔ 前端核心页面开发（Vue.js）",
    "✔ MVP Demo 可完整演示闭环流程",
]
add_bullet_textbox(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(2.5), done, font_size=15)

add_textbox(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.5), "🔭 未来规划", font_size=22, color=COLOR_PRIMARY, bold=True)
future = [
    ("短期（1-3月）", "完善核心功能，获取种子用户", COLOR_GREEN),
    ("中期（3-6月）", "扩展城市覆盖，建立合作生态", COLOR_ORANGE),
    ("长期（6-12月）", "打造宠物服务平台，多元化服务", COLOR_ACCENT),
]
for i, (period, desc, color) in enumerate(future):
    y = Inches(2.2 + i * 0.9)
    add_shape(slide, Inches(7), y, Inches(2.0), Inches(0.6), color)
    add_textbox(slide, Inches(7.1), y + Inches(0.08), Inches(1.8), Inches(0.45), period, font_size=13, color=COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(9.2), y + Inches(0.08), Inches(3.3), Inches(0.45), desc, font_size=14, color=COLOR_DARK)

add_shape(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.5), RGBColor(0xEE, 0xF0, 0xFF))
add_textbox(slide, Inches(1.2), Inches(5.3), Inches(3), Inches(0.4), "📚 团队成长", font_size=18, color=COLOR_PRIMARY, bold=True)
growth = "产品思维：学习需求分析与产品设计方法  |  技术能力：掌握全栈开发与AI模型集成  |  设计能力：提升UI/UX设计与视频剪辑技能"
add_textbox(slide, Inches(1.2), Inches(5.8), Inches(11), Inches(0.8), growth, font_size=14, color=COLOR_DARK)
add_page_number(slide, 12, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# Slide 13: 致谢
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_PRIMARY)
add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.0), "宠寻寻", font_size=48, color=COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.6), "AI图像识别寻宠平台", font_size=28, color=RGBColor(0xE0, 0xE0, 0xFF), alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(4.5), Inches(3.8), Inches(4), Inches(0.04), COLOR_WHITE)
add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.8), "让每一只走失的宠物都能安全回家", font_size=24, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5), "感谢评委老师的聆听！", font_size=22, color=RGBColor(0xCC, 0xCC, 0xFF), alignment=PP_ALIGN.CENTER)
add_page_number(slide, 13, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# Slide 14: 附录 Q&A
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLOR_WHITE)
add_page_header(slide, "附录：Q&A", "常见问题准备")

items_qa = [
    "Q1: 你们的AI识别准确率如何？",
    "    A: 采用百度AI + GPT-4 Vision双重识别，百度AI动物识别准确率超过90%，大模型可深度分析毛色、体型等特征，综合匹配准确率高。",
    "",
    "Q2: 如何保证用户隐私和数据安全？",
    "    A: 用户密码加密存储，联系方式仅在双方确认匹配后展示，敏感数据传输使用HTTPS加密。",
    "",
    "Q3: 商业模式如何实现盈利？",
    "    A: 主要通过悬赏成功后的分成、VIP增值服务、商家入驻费用和数据分析服务实现盈利。",
    "",
    "Q4: 你们的技术方案有什么优势？",
    "    A: 采用轻量级Flask + SQLite方案，部署简单成本低；AI双重识别保证准确率；Vue.js前端响应式设计，用户体验好。",
]
add_bullet_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5), items_qa, font_size=15)
add_page_number(slide, 14, TOTAL_SLIDES)

# ── 保存 ──
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "宠寻寻_答辩PPT_v2.pptx")
prs.save(output_path)
print(f'PPT generated: {output_path}')
print(f'Total slides: {len(prs.slides)}')